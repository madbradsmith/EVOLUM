#!/usr/bin/env python3
#!V_1 UPDATED FOR NEW CLEAN APP SPACE

"""
deck_builder.py

Purpose:
- Full replacement deck builder
- Removes any full-slide overlay layer that can block or dull images
- Keeps images full-bleed and visible
- Preserves localized text boxes only
- Supports legacy CLI usage from the pipeline:
    python3 deck_builder.py /home/madbrad/app/slide_plan.json
- Also supports:
    python3 deck_builder.py --project /home/madbrad/app

New in this version:
- Reads approved_brain_output.json image_plan when available
- Uses brain-directed image queries / tags for smarter stock selection
- Keeps poster / user-uploaded image priority
- Falls back to legacy folder logic when no contextual stock match exists
- Updated for renamed stock folders
"""

from __future__ import annotations

import argparse
import io
import json
import os
import re
import tempfile
import urllib.request
from pathlib import Path
from typing import Optional

from PIL import Image, ImageFilter, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

BASE_DIR = Path(__file__).resolve().parent
UPLOAD_CONTEXT_PATH = str(BASE_DIR / "user_upload_context.json")
APP_DIR = Path(__file__).resolve().parent


def load_user_context():
    try:
        with open(UPLOAD_CONTEXT_PATH, "r", encoding="utf-8") as f:
            return json.load(f)
    except:
        return {}

user_context = load_user_context()
user_poster = user_context.get("poster_filename", "")
visuals_root = APP_DIR / "visuals"
POSTER_PATH = (
    str(APP_DIR / "visuals" / "user_uploaded" / "poster" / user_poster)
    if user_poster else None
)

DEFAULT_SLIDE_PLAN = APP_DIR / "slide_plan.json"
DEFAULT_BRAIN_OUTPUT = APP_DIR / "approved_brain_output.json"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

TOP_RULE_Y = Inches(0.36)
TOP_RULE_H = Inches(0.05)

LAYOUT_THEMES = {
    "cinematic_grounded":      {"base": (18,18,22), "base2": (34,32,30), "glow": (196,126,76,44),  "accent": (206,210,219)},
    "cinematic_high_tension":  {"base": (14,10,10), "base2": (28,18,14), "glow": (220,60,40,50),   "accent": (220,100,80)},
    "contained_nocturnal":     {"base": (8,8,16),   "base2": (16,16,28), "glow": (60,80,200,40),   "accent": (100,140,220)},
    "institutional_cinematic": {"base": (12,14,18), "base2": (24,28,34), "glow": (80,120,180,36),  "accent": (160,185,210)},
    "storybook_satirical":     {"base": (20,16,10), "base2": (36,28,18), "glow": (220,180,80,44),  "accent": (220,180,80)},
    "neon_social_chaos":       {"base": (10,8,18),  "base2": (20,14,30), "glow": (180,40,240,50),  "accent": (180,80,240)},
    "athletic_prestige":       {"base": (8,14,20),  "base2": (16,26,36), "glow": (40,160,220,44),  "accent": (60,160,220)},
}

_active_theme: dict = LAYOUT_THEMES["cinematic_grounded"]


def rgb(r: int, g: int, b: int) -> RGBColor:
    return RGBColor(r, g, b)


def clean(text) -> str:
    return " ".join(str(text or "").split()).strip()


def normalize_key(text: str) -> str:
    text = clean(text).lower()
    text = re.sub(r"\s*\(.*?\)\s*", " ", text)
    text = re.sub(r"[^a-z0-9]+", " ", text)
    return " ".join(text.split())


def tokenize(text: str) -> list[str]:
    return [t for t in re.findall(r"[a-z0-9]+", clean(text).lower()) if t]


def load_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def next_output_path(output_dir: Path) -> Path:
    nums = []
    for p in output_dir.glob("pitch_deck_v*.pptx"):
        m = re.search(r"pitch_deck_v(\d+)\.pptx$", p.name)
        if m:
            nums.append(int(m.group(1)))
    nxt = max(nums) + 1 if nums else 1
    return output_dir / f"pitch_deck_v{nxt}.pptx"


def resolve_paths(args) -> tuple[Path, Path, Path]:
    if args.project:
        project_dir = Path(args.project).expanduser().resolve()
        output_dir = project_dir / "output"
        output_dir.mkdir(parents=True, exist_ok=True)
        return project_dir / "slide_plan.json", project_dir / "visuals", output_dir

    slide_plan = Path(args.slide_plan).expanduser().resolve() if args.slide_plan else DEFAULT_SLIDE_PLAN.resolve()
    project_dir = slide_plan.parent
    output_dir = project_dir / "output"
    output_dir.mkdir(parents=True, exist_ok=True)
    return slide_plan, project_dir / "visuals", output_dir


def load_brain_output(project_dir: Path) -> dict:
    candidate = project_dir / "approved_brain_output.json"
    if candidate.exists():
        return load_json(candidate)
    if DEFAULT_BRAIN_OUTPUT.exists():
        return load_json(DEFAULT_BRAIN_OUTPUT)
    return {}


_stock_image_files_cache: dict[str, list[Path]] = {}

def _stock_image_files(visuals_dir: Path, exts: set[str]) -> list[Path]:
    cache_key = str(visuals_dir)
    if cache_key in _stock_image_files_cache:
        return _stock_image_files_cache[cache_key]

    if not visuals_dir.exists():
        return []

    import re

    numbered_top_dirs = []

    for child in visuals_dir.iterdir():
        if child.is_dir() and re.match(r"^\d{2}_", child.name):
            numbered_top_dirs.append(child)

    numbered_top_dirs.sort(key=lambda p: p.name.lower())

    files = []
    for top_dir in numbered_top_dirs:
        for p in top_dir.rglob("*"):
            if p.is_file() and p.suffix.lower() in exts:
                files.append(p)

    _stock_image_files_cache[cache_key] = files
    return files

_stock_rotation_counters: dict[str, int] = {}
_user_rotation_counters: dict[str, int] = {}
_brain_stock_rotation_counters: dict[str, int] = {}
_brain_folder_usage_counts: dict[str, int] = {}
_image_usage_counts: dict[str, int] = {}


def _image_usage_key(path: Path) -> str:
    try:
        return str(path.relative_to(APP_DIR / "visuals")).lower()
    except Exception:
        return str(path).lower()


def _image_use_count(path: Path) -> int:
    return _image_usage_counts.get(_image_usage_key(path), 0)


def _mark_image_used(path: Optional[Path]) -> None:
    if not path:
        return
    key = _image_usage_key(path)
    _image_usage_counts[key] = _image_usage_counts.get(key, 0) + 1


def _pick_candidate_with_repeat_control(candidates: list[Path], start_idx: int, last_used_name: str = "") -> tuple[Optional[Path], int]:
    if not candidates:
        return None, start_idx

    ordered = [candidates[(start_idx + offset) % len(candidates)] for offset in range(len(candidates))]

    # Pass 1: best case — brand new image, never adjacent repeat.
    for offset, candidate in enumerate(ordered):
        if candidate.name == last_used_name:
            continue
        if _image_use_count(candidate) == 0:
            return candidate, (start_idx + offset + 1) % len(candidates)

    # Pass 2: allow one reuse, but never more than twice in a deck.
    for offset, candidate in enumerate(ordered):
        if candidate.name == last_used_name:
            continue
        if _image_use_count(candidate) == 1:
            return candidate, (start_idx + offset + 1) % len(candidates)

    # Pass 3: emergency fallback — still never repeat adjacent slides,
    # and still block anything already used twice or more.
    for offset, candidate in enumerate(ordered):
        if candidate.name == last_used_name:
            continue
        if _image_use_count(candidate) < 2:
            return candidate, (start_idx + offset + 1) % len(candidates)

    # If the image pool is truly exhausted, return None so the caller can
    # fall through to a different selection source instead of repeating forever.
    return None, start_idx


def reset_image_selection_state() -> None:
    _stock_rotation_counters.clear()
    _user_rotation_counters.clear()
    _brain_stock_rotation_counters.clear()
    _brain_folder_usage_counts.clear()
    _image_usage_counts.clear()


def _top_visual_folder(path: Path) -> str:
    try:
        rel_parts = path.relative_to(APP_DIR / "visuals").parts
        return rel_parts[0] if rel_parts else ""
    except Exception:
        return path.parent.name


def _select_user_image(current_files: list[Path], lookup_key: str, last_used_name: str = "") -> Optional[Path]:
    if not current_files:
        return None, "none"

    candidates = sorted(current_files, key=lambda p: (
        _image_use_count(p),
        p.name.lower(),
    ))
    start_idx = _user_rotation_counters.get(lookup_key, 0)
    candidate, next_idx = _pick_candidate_with_repeat_control(candidates, start_idx, last_used_name)
    _user_rotation_counters[lookup_key] = next_idx
    return candidate


def _stock_candidates_for_key(stock_files: list[Path], lookup_key: str) -> list[Path]:
    folder_map = {
        "__title__": ["01_cinematic_tension", "03_urban_pressure", "07_night_isolation"],
        "logline": ["03_urban_pressure", "07_night_isolation", "01_cinematic_tension"],
        "synopsis": ["02_emotional_grounded", "03_urban_pressure", "07_night_isolation"],
        "protagonist": ["02_emotional_grounded", "06_controlled_clean", "03_urban_pressure"],
        "antagonist": ["03_urban_pressure", "01_cinematic_tension", "07_night_isolation"],
        "supporting characters": ["02_emotional_grounded", "03_urban_pressure", "06_controlled_clean"],
        "theme": ["02_emotional_grounded", "06_controlled_clean", "08_daylight_release"],
        "tone": ["01_cinematic_tension", "07_night_isolation", "02_emotional_grounded"],
        "world": ["03_urban_pressure", "08_daylight_release", "07_night_isolation"],
        "conflict engine": ["03_urban_pressure", "07_night_isolation", "01_cinematic_tension"],
        "stakes": ["07_night_isolation", "03_urban_pressure", "01_cinematic_tension"],
        "why this film": ["01_cinematic_tension", "03_urban_pressure", "06_controlled_clean"],
        "audience": ["02_emotional_grounded", "06_controlled_clean", "08_daylight_release"],
        "visual style": ["01_cinematic_tension", "07_night_isolation", "02_emotional_grounded"],
        "comparables": ["01_cinematic_tension", "06_controlled_clean", "03_urban_pressure"],
        "market position": ["06_controlled_clean", "01_cinematic_tension", "03_urban_pressure"],
        "market projections": ["06_controlled_clean", "01_cinematic_tension", "09_institutional_authority"],
        "director vision": ["01_cinematic_tension", "02_emotional_grounded", "03_urban_pressure"],
        "casting ideas": ["02_emotional_grounded", "06_controlled_clean", "03_urban_pressure"],
        "production scope": ["06_controlled_clean", "08_daylight_release", "02_emotional_grounded"],
        "closing statement": ["01_cinematic_tension", "07_night_isolation", "02_emotional_grounded"],

        # legacy keys
        "hook": ["07_night_isolation", "03_urban_pressure", "01_cinematic_tension"],
        "conflict": ["03_urban_pressure", "07_night_isolation", "01_cinematic_tension"],
        "story engine": ["03_urban_pressure", "07_night_isolation", "06_controlled_clean"],
        "reversal": ["01_cinematic_tension", "07_night_isolation", "03_urban_pressure"],
        "themes": ["02_emotional_grounded", "06_controlled_clean", "08_daylight_release"],
        "why this movie": ["01_cinematic_tension", "03_urban_pressure", "06_controlled_clean"],
    }

    mapped_folders = folder_map.get(lookup_key, folder_map["synopsis"])
    by_folder: dict[str, list[Path]] = {folder: [] for folder in mapped_folders}

    for p in stock_files:
        try:
            rel_parts = p.relative_to(visuals_root).parts if str(visuals_root) in str(p) else p.parts
            top = rel_parts[0] if rel_parts else ""
            if top in by_folder:
                by_folder[top].append(p)
        except Exception:
            continue

    candidates: list[Path] = []
    for folder in mapped_folders:
        folder_files = sorted(by_folder.get(folder, []), key=lambda x: x.name.lower())
        candidates.extend(folder_files)

    return candidates


def _select_stock_image(stock_files: list[Path], lookup_key: str, last_used_name: str = "") -> Optional[Path]:
    candidates = _stock_candidates_for_key(stock_files, lookup_key)
    if not candidates:
        return None

    candidates = sorted(candidates, key=lambda p: (
        _image_use_count(p),
        p.name.lower(),
    ))

    key = lookup_key
    start_idx = _stock_rotation_counters.get(key, 0)
    candidate, next_idx = _pick_candidate_with_repeat_control(candidates, start_idx, last_used_name)
    _stock_rotation_counters[key] = next_idx
    return candidate


def _brain_image_instruction(brain_output: dict, slide_title: str, slide_number: int) -> Optional[dict]:
    image_plan = brain_output.get("image_plan", [])
    if not isinstance(image_plan, list):
        return None

    normalized_title = normalize_key(slide_title)

    for item in image_plan:
        if not isinstance(item, dict):
            continue
        if item.get("slide_number") == slide_number:
            return item

    for item in image_plan:
        if not isinstance(item, dict):
            continue
        if normalize_key(item.get("slide_title", "")) == normalized_title:
            return item

    if slide_number == 1:
        for item in image_plan:
            if not isinstance(item, dict):
                continue
            if normalize_key(item.get("slide_title", "")) == "title":
                return item

    return None


def _score_stock_file_against_tags(path: Path, tags: list[str]) -> int:
    try:
        rel_text = normalize_key(str(path.relative_to(APP_DIR / "visuals")))
    except Exception:
        rel_text = normalize_key(str(path))
    filename_text = normalize_key(path.stem)
    combined = f"{rel_text} {filename_text}"

    score = 0
    for tag in tags:
        tag_norm = normalize_key(tag)
        if not tag_norm:
            continue
        parts = tag_norm.split()

        if tag_norm in combined:
            score += 8

        matched_parts = sum(1 for part in parts if part in combined)
        score += matched_parts * 3

    return score


def _select_brain_directed_stock_image(
    stock_files: list[Path],
    image_instruction: Optional[dict],
    slide_title: str,
    last_used_name: str = ""
) -> Optional[Path]:
    if not image_instruction:
        return None

    tags = image_instruction.get("image_tags") or []
    query = image_instruction.get("image_query", "")

    if isinstance(query, str) and query.strip():
        tags = list(tags) + query.split()

    normalized_tags = []
    seen = set()
    for tag in tags:
        tag_norm = normalize_key(str(tag))
        if tag_norm and tag_norm not in seen:
            seen.add(tag_norm)
            normalized_tags.append(tag_norm)

    if not normalized_tags:
        return None

    scored: list[tuple[int, Path]] = []
    for p, combined in _get_precomputed(stock_files):
        score = _score_combined(combined, normalized_tags)
        if score > 0:
            scored.append((score, p))

    if not scored:
        return None

    # Keep only close-score candidates so we preserve relevance,
    # then prefer folders that have been used less in this build.
    best_raw_score = max(score for score, _ in scored)
    score_window = 6
    close_scored = [(score, p) for score, p in scored if score >= best_raw_score - score_window]

    close_scored.sort(
        key=lambda item: (
            _image_use_count(item[1]),
            _brain_folder_usage_counts.get(_top_visual_folder(item[1]), 0),
            -item[0],
            item[1].name.lower(),
        )
    )
    candidates = [p for _, p in close_scored]

    rotation_key = f"brain::{normalize_key(slide_title)}"
    start_idx = _brain_stock_rotation_counters.get(rotation_key, 0)
    candidate, next_idx = _pick_candidate_with_repeat_control(candidates, start_idx, last_used_name)
    if not candidate:
        return None

    _brain_stock_rotation_counters[rotation_key] = next_idx
    folder_key = _top_visual_folder(candidate)
    _brain_folder_usage_counts[folder_key] = _brain_folder_usage_counts.get(folder_key, 0) + 1
    return candidate


def resolve_image_options_for_slide(
    visuals_dir: Path,
    slide_info: dict,
    image_for_slide: Optional[Path],
    image_source: str,
    slide_title: str,
) -> list[dict]:
    exts = {".png", ".jpg", ".jpeg", ".webp", ".PNG", ".JPG", ".JPEG", ".WEBP"}
    stock_files = _stock_image_files(visuals_dir, exts)
    raw_options = slide_info.get("image_options") or []
    resolved_options: list[dict] = []
    seen_paths: set[str] = set()

    def add_option(payload: dict) -> None:
        path_val = str(payload.get("image_path", "") or "").strip()
        if not path_val or path_val in seen_paths:
            return
        seen_paths.add(path_val)
        resolved_options.append(payload)

    if image_for_slide:
        add_option({
            "rank": 1,
            "option_id": "selected",
            "label": "Current Pick",
            "focus": "selected",
            "image_path": str(image_for_slide),
            "image_name": image_for_slide.name,
            "image_source": image_source,
        })

    if isinstance(raw_options, list):
        for option in raw_options:
            if not isinstance(option, dict):
                continue
            instruction = {
                "image_query": option.get("image_query", ""),
                "image_tags": option.get("image_tags", []),
            }
            option_rank = int(option.get("rank", len(resolved_options) + 1) or (len(resolved_options) + 1))
            option_id = clean(option.get("option_id") or f"option_{option_rank}")
            option_label = clean(option.get("label") or f"Option {option_rank}")
            option_focus = clean(option.get("focus") or "alternate")

            resolved = _select_brain_directed_stock_image(
                stock_files,
                instruction,
                f"{slide_title}_{option_id}",
                last_used_name="",
            )
            if not resolved:
                continue

            add_option({
                "rank": option_rank,
                "option_id": option_id,
                "label": option_label,
                "focus": option_focus,
                "image_path": str(resolved),
                "image_name": resolved.name,
                "image_source": "brain_stock_option",
            })

            if len(resolved_options) >= 4:
                break

    # Fill remaining slots with images from unused folders for variety
    if len(resolved_options) < 4 and stock_files:
        used_paths = {o["image_path"] for o in resolved_options}
        used_folders = {str(Path(p).parent) for p in used_paths}

        by_folder: dict[str, list[Path]] = {}
        for f in stock_files:
            folder = str(f.parent)
            by_folder.setdefault(folder, []).append(f)

        fill_rank = len(resolved_options) + 1
        for folder, files in sorted(by_folder.items()):
            if len(resolved_options) >= 4:
                break
            if folder in used_folders:
                continue
            candidate = next((f for f in files if str(f) not in used_paths), None)
            if not candidate:
                continue
            used_paths.add(str(candidate))
            used_folders.add(folder)
            add_option({
                "rank": fill_rank,
                "option_id": f"fill_{fill_rank}",
                "label": f"Alt {fill_rank - 1}",
                "focus": "alternate",
                "image_path": str(candidate),
                "image_name": candidate.name,
                "image_source": "folder_fill",
            })
            fill_rank += 1

    for idx, option in enumerate(resolved_options, start=1):
        option["rank"] = idx

    return resolved_options[:5]


FAL_API_KEY = os.environ.get("FAL_API_KEY", "")
EVOLUM_SESSION_ID = os.environ.get("EVOLUM_SESSION_ID", "shared")

_SLIDE_VISUAL_CONCEPTS = {
    "logline":              "cinematic establishing shot, wide angle, dramatic lighting",
    "synopsis":             "cinematic scene, atmospheric, narrative moment",
    "synopsis 2":           "cinematic scene, mid-shot, dramatic tension",
    "synopsis 3":           "cinematic scene, close-up, emotional intensity",
    "protagonist":          "cinematic portrait, single character, dramatic lighting, film still",
    "antagonist":           "cinematic portrait, menacing figure, dramatic shadows, film still",
    "supporting characters":"cinematic ensemble shot, multiple characters, film still",
    "world":                "cinematic landscape, establishing shot, rich environment",
    "hook":                 "cinematic close-up, tension, dramatic moment",
    "conflict":             "cinematic confrontation, dramatic tension, high stakes",
    "stakes":               "cinematic wide shot, weight of consequence, dramatic",
    "tone":                 "cinematic mood shot, atmospheric lighting, visual tone",
    "story engine":         "cinematic action, driving force, momentum",
    "reversal":             "cinematic turning point, dramatic shift, pivotal moment",
    "themes":               "cinematic symbolic imagery, thematic visual metaphor",
    "why this movie":       "cinematic wide shot, cultural moment, compelling imagery",
    "comparables":          "cinematic collage feel, prestige film aesthetic",
    "market projections":   "cinematic wide shot, commercial appeal, high production value",
    "closing statement":    "cinematic final frame, powerful, memorable",
    "closing":              "cinematic hero shot, wide epic frame, powerful final image, golden light",
}

_GENRE_STYLE = {
    "horror":       "dark, unsettling, atmospheric horror, shadows, practical effects aesthetic",
    "thriller":     "tense, noir-influenced, sharp contrast, suspenseful",
    "comedy":       "warm lighting, vibrant colors, playful composition",
    "drama":        "naturalistic lighting, intimate, emotionally grounded",
    "action":       "dynamic, kinetic energy, bold framing, high contrast",
    "sci-fi":       "futuristic, cool tones, technological, epic scale",
    "fantasy":      "magical, rich colors, otherworldly, painterly lighting",
    "romance":      "warm golden tones, soft focus, intimate, emotional",
    "documentary":  "gritty realism, candid, natural light, observational",
    "animation":    "stylized, vibrant, expressive, dynamic",
}


_PERIOD_KEYWORDS = [
    (["medieval", "kingdom", "castle", "court", "jester", "knight", "throne", "king", "queen", "lord", "valoria"], "medieval fantasy, period accurate, castle, kingdom aesthetic"),
    (["space", "spaceship", "galaxy", "planet", "alien", "starship", "orbit"], "outer space, sci-fi, futuristic spacecraft"),
    (["future", "cyberpunk", "neon", "dystopia", "android", "robot"], "near-future dystopia, cyberpunk neon lighting"),
    (["western", "frontier", "cowboy", "saloon", "sheriff"], "american western, frontier, dusty plains"),
    (["victorian", "1800s", "19th century", "gaslight", "corset"], "victorian era, 1800s period costume"),
    (["1920", "1930", "prohibition", "jazz age", "noir", "gangster"], "1930s noir, prohibition era, art deco"),
    (["war", "wwii", "battlefield", "soldier", "trench", "military"], "wartime, gritty military realism"),
    (["ancient", "roman", "greek", "egypt", "pyramid", "colosseum"], "ancient world, epic historical"),
    (["pirate", "ship", "ocean", "sail", "treasure", "sea"], "age of sail, pirate adventure, tall ships"),
    (["animation", "animated", "cartoon", "pixar", "anime"], "stylized animation, vibrant illustration"),
]


def _detect_period_style(world: str, genre: str, tone: str) -> str:
    combined = f"{world} {genre} {tone}".lower()
    for keywords, style in _PERIOD_KEYWORDS:
        if any(k in combined for k in keywords):
            return style
    return ""


def build_image_prompt(slide_title: str, brain_output: dict) -> str:
    normalized = normalize_key(slide_title)
    concept = _SLIDE_VISUAL_CONCEPTS.get(normalized, "cinematic scene, dramatic lighting, film still")

    genre = str(brain_output.get("genre", "drama")).lower()
    genre_style = ""
    for g, style in _GENRE_STYLE.items():
        if g in genre:
            genre_style = style
            break
    if not genre_style:
        genre_style = "cinematic, naturalistic lighting, film aesthetic"

    tone = str(brain_output.get("tone", "")).lower()
    world = str(brain_output.get("world", "")).replace("\n", " ").strip()

    period_style = _detect_period_style(world, genre, tone)
    period_hint = f", {period_style}" if period_style else ""
    world_hint = f", {world[:80]}" if world and not period_style else ""
    tone_hint = f", {tone[:60]}" if tone else ""

    prompt = (
        f"{concept}, {genre_style}{period_hint}{world_hint}{tone_hint}, "
        f"professional film still, 35mm, shallow depth of field, no text, no watermarks, "
        f"ultra-detailed, photorealistic, 16:9 aspect ratio"
    )
    return prompt


def generate_fal_image(prompt: str, cache_path: Path) -> Optional[Path]:
    if not FAL_API_KEY:
        return None
    if cache_path.exists():
        return cache_path

    import urllib.error
    url = "https://fal.run/fal-ai/flux/schnell"
    payload = json.dumps({
        "prompt": prompt,
        "image_size": "landscape_16_9",
        "num_inference_steps": 4,
        "num_images": 1,
        "enable_safety_checker": True,
    }).encode("utf-8")

    req = urllib.request.Request(
        url,
        data=payload,
        headers={
            "Authorization": f"Key {FAL_API_KEY}",
            "Content-Type": "application/json",
        },
        method="POST",
    )
    try:
        with urllib.request.urlopen(req, timeout=60) as resp:
            result = json.loads(resp.read().decode("utf-8"))
        image_url = result["images"][0]["url"]
        cache_path.parent.mkdir(parents=True, exist_ok=True)
        urllib.request.urlretrieve(image_url, cache_path)
        print(f"✨ FAL generated image for prompt: {prompt[:60]}...")
        return cache_path
    except Exception as e:
        print(f"⚠️  FAL image generation failed: {e}")
        return None


def find_image_for_slide(
    visuals_dir: Path,
    deck_title: str,
    slide_title: str,
    slide_number: int,
    brain_output: Optional[dict] = None,
    last_used_name: str = ""
) -> tuple[Optional[Path], str]:
    poster_dir = visuals_dir / "user_uploaded" / "poster"
    current_dir = visuals_dir / "user_uploaded" / "current"
    exts = {".png", ".jpg", ".jpeg", ".webp", ".PNG", ".JPG", ".JPEG", ".WEBP"}

    poster_files = [p for p in poster_dir.glob("*") if p.suffix in exts] if poster_dir.exists() else []
    current_files = [p for p in current_dir.glob("*") if p.suffix in exts] if current_dir.exists() else []
    stock_files = _stock_image_files(visuals_dir, exts)

    normalized_title = normalize_key(slide_title)
    image_instruction = _brain_image_instruction(brain_output or {}, slide_title, slide_number)

    non_title_slides = {
        "logline", "synopsis", "synopsis 2", "synopsis 3", "synopsis 4",
        "protagonist", "antagonist", "supporting characters", "world", "hook",
        "conflict", "conflict engine", "stakes", "tone", "story engine",
        "reversal", "theme", "themes", "why this movie", "why this film",
        "audience", "visual style", "comparables", "market position",
        "market projections", "director vision", "casting ideas", "production scope",
        "closing statement", "closing"
    }

    if normalized_title not in non_title_slides:
        for p in poster_files:
            print(f"🖼️ Using POSTER for '{deck_title}': {p}")
            return p, "poster"

        brain_title = _select_brain_directed_stock_image(stock_files, image_instruction, slide_title, last_used_name)
        if brain_title:
            print(f"🖼️ Using BRAIN STOCK image for title '{deck_title}': {brain_title}")
            return brain_title, "brain_stock"

        stock_title = _select_stock_image(stock_files, "__title__", last_used_name)
        if stock_title:
            print(f"🖼️ Using STOCK image for title '{deck_title}': {stock_title}")
            return stock_title, "stock"

    key_map = {
        "logline": ["SceneD", "Misdirect", "logline", "frame_1"],
        "synopsis": ["SceneE", "Decision", "rear_pov", "whisper", "synopsis", "frame_2"],
        "protagonist": ["malik", "character", "worldslide", "frame_3"],
        "world": ["world", "city", "character_slide", "frame_4"],
        "hook": ["look_up", "hook", "frame_5"],
        "conflict": ["rear_pov", "conflict", "frame_6"],
        "stakes": ["whisper", "stakes"],
        "tone": ["rachel", "tone"],
        "story engine": ["malik", "engine"],
        "reversal": ["Decision", "reversal", "SceneE"],
        "themes": ["rachel", "themes"],
        "why this movie": ["rico", "why"],
    }

    lookup_key = "synopsis" if normalized_title.startswith("synopsis") else normalized_title

    if lookup_key in key_map:
        matched_files: list[Path] = []
        for needle in key_map[lookup_key]:
            for p in current_files:
                if needle.lower() in p.name.lower():
                    if p not in matched_files:
                        matched_files.append(p)

        if matched_files:
            selected = _select_user_image(matched_files, lookup_key, last_used_name)
            if selected:
                print(f"🖼️ Using USER image for '{slide_title}': {selected}")
                return selected, "user"

    if current_files:
        selected = _select_user_image(current_files, lookup_key, last_used_name)
        if selected:
            print(f"🖼️ Using USER image fallback for '{slide_title}': {selected}")
            return selected, "user_fallback"

    if FAL_API_KEY and brain_output:
        prompt = build_image_prompt(slide_title, brain_output)
        cache_dir = visuals_dir.parent / "generated_images" / EVOLUM_SESSION_ID
        safe_title = re.sub(r"[^a-z0-9_]", "_", normalized_title)
        cache_path = cache_dir / f"{slide_number:02d}_{safe_title}.jpg"
        generated = generate_fal_image(prompt, cache_path)
        if generated:
            return generated, "fal_generated"

    brain_stock = _select_brain_directed_stock_image(stock_files, image_instruction, slide_title, last_used_name)
    if brain_stock:
        print(f"🖼️ Using BRAIN STOCK image for '{slide_title}': {brain_stock}")
        return brain_stock, "brain_stock"

    stock_match = _select_stock_image(stock_files, lookup_key, last_used_name)
    if stock_match:
        print(f"🖼️ Using STOCK image for '{slide_title}': {stock_match}")
        return stock_match, "stock"

    return None, "none"


# Pre-computed (path, combined_search_text) pairs — keyed by id of the cached file list
_stock_precomputed_cache: dict[int, list[tuple]] = {}


def _get_precomputed(stock_files: list[Path]) -> list[tuple]:
    """Return (path, combined_normalized_text) for every file — computed once per list object."""
    cache_key = id(stock_files)
    if cache_key not in _stock_precomputed_cache:
        visuals_base = APP_DIR / "visuals"
        result = []
        for p in stock_files:
            try:
                rel_text = normalize_key(str(p.relative_to(visuals_base)))
            except Exception:
                rel_text = normalize_key(str(p))
            combined = f"{rel_text} {normalize_key(p.stem)}"
            result.append((p, combined))
        _stock_precomputed_cache[cache_key] = result
    return _stock_precomputed_cache[cache_key]


def _score_combined(combined: str, normalized_tags: list[str]) -> int:
    score = 0
    for tag in normalized_tags:
        if not tag:
            continue
        if tag in combined:
            score += 8
        score += sum(3 for part in tag.split() if part in combined)
    return score


# Cache rendered base backgrounds — one expensive render per theme per process
_base_bg_cache: dict[int, bytes] = {}


def _render_base_bg(theme: dict) -> bytes:
    width_px, height_px = 640, 360
    b1 = theme["base"]
    b2 = theme["base2"]
    glow_color = theme["glow"]

    img = Image.new("RGB", (width_px, height_px), b1)
    draw = ImageDraw.Draw(img)
    for y in range(height_px):
        t = y / max(1, height_px - 1)
        r = int(b1[0] + (b2[0] - b1[0]) * t)
        g = int(b1[1] + (b2[1] - b1[1]) * t)
        b = int(b1[2] + (b2[2] - b1[2]) * t)
        draw.line((0, y, width_px, y), fill=(r, g, b))

    glow = Image.new("RGBA", (width_px, height_px), (0, 0, 0, 0))
    gdraw = ImageDraw.Draw(glow)
    gdraw.ellipse(
        (width_px * 0.18, height_px * 0.15, width_px * 0.82, height_px * 0.95),
        fill=glow_color,
    )
    glow = glow.filter(ImageFilter.GaussianBlur(40))
    img = Image.alpha_composite(img.convert("RGBA"), glow).convert("RGB")

    vignette = Image.new("L", (width_px, height_px), 0)
    vdraw = ImageDraw.Draw(vignette)
    vdraw.ellipse(
        (-width_px * 0.05, -height_px * 0.12, width_px * 1.05, height_px * 1.08),
        fill=220,
    )
    vignette = vignette.filter(ImageFilter.GaussianBlur(50))
    dark = Image.new("RGB", (width_px, height_px), (10, 10, 12))
    img = Image.composite(img, dark, vignette)

    buf = io.BytesIO()
    img.save(buf, format="JPEG", quality=80, optimize=True)
    return buf.getvalue()


def add_base_background(slide) -> None:
    theme = _active_theme
    cache_key = id(theme)
    if cache_key not in _base_bg_cache:
        _base_bg_cache[cache_key] = _render_base_bg(theme)

    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".jpg")
    tmp.write(_base_bg_cache[cache_key])
    tmp.close()
    slide.shapes.add_picture(tmp.name, 0, 0, width=SLIDE_W, height=SLIDE_H)
    os.unlink(tmp.name)


def add_blur_background(slide, image_path: Optional[Path]) -> None:
    if not image_path or not image_path.exists():
        return

    with Image.open(image_path) as im:
        bg = im.convert("RGB")
        bg.thumbnail((320, 180))
        bg = bg.resize((640, 360))
        bg = bg.filter(ImageFilter.GaussianBlur(8))
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".jpg")
        bg.save(tmp.name, format="JPEG", quality=65, optimize=True)

    slide.shapes.add_picture(str(tmp.name), 0, 0, width=SLIDE_W, height=SLIDE_H)
    import os; os.unlink(tmp.name)


def add_center_image(slide, image_path: Optional[Path], scale_factor: float = 0.68) -> None:
    if not image_path or not image_path.exists():
        return

    with Image.open(image_path) as im:
        img = im.convert("RGB")
        img.thumbnail((960, 540))
        img_w, img_h = img.size

        if img_w <= 0 or img_h <= 0:
            return

        scale = min(float(SLIDE_W) / img_w, float(SLIDE_H) / img_h) * scale_factor
        render_w = int(img_w * scale)
        render_h = int(img_h * scale)
        left = int((float(SLIDE_W) - render_w) / 2)
        top = int((float(SLIDE_H) - render_h) / 2)

        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".jpg")
        img.save(tmp.name, format="JPEG", quality=78, optimize=True)

    slide.shapes.add_picture(str(tmp.name), left, top, width=render_w, height=render_h)
    import os; os.unlink(tmp.name)


def add_full_bleed_image(slide, image_path: Optional[Path]) -> None:
    add_blur_background(slide, image_path)
    add_center_image(slide, image_path, scale_factor=0.68)


def add_title_poster_image(slide, image_path: Optional[Path]) -> None:
    add_blur_background(slide, image_path)
    add_center_image(slide, image_path, scale_factor=0.82)


def add_top_rule(slide) -> None:
    return


def add_title_text(slide, text: str) -> None:
    tx = slide.shapes.add_textbox(Inches(0.7), Inches(0.58), Inches(12.0), Inches(0.45))
    tf = tx.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = clean(text)
    run.font.size = Pt(17)
    run.font.bold = True
    run.font.color.rgb = rgb(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER


def add_text_box(slide, left, top, width, height, text: str, *, font_size: int = 18,
                 align=PP_ALIGN.LEFT, fill_transparency: float = 0.22) -> None:
    accent = _active_theme["accent"]
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = rgb(14, 14, 18)
    box.fill.transparency = fill_transparency
    box.line.color.rgb = rgb(*accent)
    box.line.width = Pt(1.2)

    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.24)
    tf.margin_right = Inches(0.24)
    tf.margin_top = Inches(0.16)
    tf.margin_bottom = Inches(0.12)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = clean(text)
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = rgb(255, 255, 255)
    p.alignment = align


def add_cinematic_caption(slide, body: str, font_size: int = 18) -> None:
    """Full-width dark band anchored at the bottom — no border, text sits on the image."""
    if not body:
        return
    band_h = Inches(1.55)
    band_top = SLIDE_H - band_h
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, band_top, SLIDE_W, band_h)
    band.fill.solid()
    band.fill.fore_color.rgb = rgb(6, 6, 8)
    band.fill.transparency = 0.18
    band.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(0.6), band_top, SLIDE_W - Inches(1.2), band_h)
    tf = tx.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.18)
    tf.margin_bottom = Inches(0.1)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = clean(body)
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = rgb(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER


def _auto_font_size(text: str, base: int) -> int:
    n = len(text)
    if n > 300: return max(10, base - 5)
    if n > 220: return max(11, base - 4)
    if n > 160: return max(12, base - 3)
    if n > 110: return max(13, base - 2)
    if n > 70:  return max(14, base - 1)
    return base


def build_slide_split_panel(slide, image_path: Optional[Path], slide_title: str, body: str) -> None:
    """Layout B — image fills left 55%, dark text panel right 45%."""
    add_base_background(slide)
    accent = _active_theme["accent"]

    panel_w = int(float(SLIDE_W) * 0.55)
    panel_h = int(float(SLIDE_H))

    # Pixel dimensions for image processing (EMU units above are not pixels)
    panel_w_px = 704
    panel_h_px = 720

    if image_path and image_path.exists():
        try:
            with Image.open(image_path) as im:
                img = im.convert("RGB")
                img.thumbnail((panel_w_px * 2, panel_h_px * 2))
                img_ratio = img.width / img.height
                panel_ratio = panel_w_px / panel_h_px
                if img_ratio > panel_ratio:
                    new_h = panel_h_px
                    new_w = int(new_h * img_ratio)
                else:
                    new_w = panel_w_px
                    new_h = int(new_w / img_ratio)
                img = img.resize((new_w, new_h), Image.LANCZOS)
                lc = (new_w - panel_w_px) // 2
                tc = (new_h - panel_h_px) // 2
                img = img.crop((lc, tc, lc + panel_w_px, tc + panel_h_px))
                tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".jpg")
                img.save(tmp.name, format="JPEG", quality=82, optimize=True)
            slide.shapes.add_picture(str(tmp.name), 0, 0, width=panel_w, height=SLIDE_H)
            os.unlink(tmp.name)
        except Exception:
            pass

    # Dark right panel
    right_x = int(float(SLIDE_W) * 0.54)
    right_w = int(float(SLIDE_W) - right_x)
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, right_x, 0, right_w, SLIDE_H)
    panel.fill.solid()
    panel.fill.fore_color.rgb = rgb(10, 10, 14)
    panel.fill.transparency = 0.0
    panel.line.fill.background()

    # Accent divider line
    div = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, right_x, 0, Inches(0.04), SLIDE_H)
    div.fill.solid()
    div.fill.fore_color.rgb = rgb(*accent)
    div.fill.transparency = 0.35
    div.line.fill.background()

    # Title in right panel
    tx_left = right_x + int(Inches(0.28))
    tx_w = right_w - int(Inches(0.56))
    tx = slide.shapes.add_textbox(tx_left, Inches(0.48), tx_w, Inches(0.9))
    tf = tx.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = clean(slide_title.split("(")[0].strip())
    run.font.size = Pt(15)
    run.font.bold = True
    run.font.color.rgb = rgb(*accent)
    p.alignment = PP_ALIGN.LEFT

    # Body in right panel — generous height, auto font
    font_size = _auto_font_size(body, base=17)
    add_text_box(slide, tx_left, Inches(1.55), tx_w, Inches(5.2),
                 body, font_size=font_size, align=PP_ALIGN.LEFT, fill_transparency=0.0)


def build_slide_text_only(slide, slide_title: str, body: str) -> None:
    """Text-only layout — no image. Big centered body text fills the slide."""
    add_base_background(slide)
    accent = _active_theme["accent"]
    add_top_rule(slide)

    # Title band
    tx = slide.shapes.add_textbox(Inches(1.0), Inches(0.52), Inches(11.3), Inches(0.7))
    tf = tx.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = clean(slide_title.split("(")[0].strip())
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = rgb(*accent)
    p.alignment = PP_ALIGN.CENTER

    # Large body — fills most of slide
    font_size = _auto_font_size(body, base=22)
    font_size = max(font_size, 20)
    add_text_box(slide, Inches(1.0), Inches(1.5), Inches(11.3), Inches(5.4),
                 body, font_size=font_size, align=PP_ALIGN.CENTER, fill_transparency=0.0)


def build_slide_editorial(slide, image_path: Optional[Path], slide_title: str, body: str) -> None:
    """Layout C — image floats center-top, title below it, wide body box at bottom."""
    add_base_background(slide)
    accent = _active_theme["accent"]

    if image_path and image_path.exists():
        add_center_image(slide, image_path, scale_factor=0.52)

    # Title centered below image area
    tx = slide.shapes.add_textbox(Inches(1.0), Inches(4.3), Inches(11.3), Inches(0.6))
    tf = tx.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = clean(slide_title.split("(")[0].strip())
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = rgb(*accent)
    p.alignment = PP_ALIGN.CENTER

    # Wide body box at bottom
    font_size = _auto_font_size(body, base=16)
    add_text_box(slide, Inches(0.7), Inches(5.1), Inches(11.93), Inches(2.0),
                 body, font_size=font_size, align=PP_ALIGN.CENTER, fill_transparency=0.18)


def place_text_by_stage(slide, stage: str, layout: str, body: str) -> None:
    stage = clean(stage).lower()
    layout = clean(layout).lower()
    if not body:
        return
    fs = _auto_font_size(body, 18)
    if layout in {"title", "closing"}:
        fs = _auto_font_size(body, 19)
    add_cinematic_caption(slide, body, font_size=fs)


def build_presentation(slide_plan_path: Path, visuals_dir: Path, output_dir: Path) -> Path:
    global _active_theme
    reset_image_selection_state()
    plan = load_json(slide_plan_path)
    brain_output = load_brain_output(output_dir)

    layout_strategy = brain_output.get("layout_strategy") or {}
    layout_style = (layout_strategy.get("layout_style") or "cinematic_grounded").strip()
    composition_bias = (layout_strategy.get("composition_bias") or "image_forward").strip()
    _active_theme = LAYOUT_THEMES.get(layout_style, LAYOUT_THEMES["cinematic_grounded"])
    print(f"🎨 Layout theme: {layout_style} | Composition: {composition_bias}")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    deck_title = clean(plan.get("title", "Project"))
    manifest: list[dict] = []

    last_used_image_name = ""

    for idx, slide_info in enumerate(plan.get("slides", []), start=1):
        slide_title = clean(slide_info.get("title"))
        body = clean(slide_info.get("body"))
        layout = clean(slide_info.get("layout"))
        stage = clean(slide_info.get("stage"))
        slide_number = int(slide_info.get("slide_number", idx))

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        explicit_path_str = str(slide_info.get("image_path") or "").strip()
        if explicit_path_str == "__none__":
            image_for_slide = None
            image_source = "text_only"
        elif explicit_path_str:
            explicit = Path(explicit_path_str)
            if not explicit.is_absolute():
                explicit = (APP_DIR / explicit).resolve()
            if explicit.exists():
                image_for_slide = explicit
                image_source = str(slide_info.get("image_source") or "user_selected")
            else:
                image_for_slide, image_source = find_image_for_slide(
                    visuals_dir=visuals_dir,
                    deck_title=deck_title,
                    slide_title=slide_title if layout != "title" else deck_title,
                    slide_number=slide_number,
                    brain_output=brain_output,
                    last_used_name=last_used_image_name
                )
        else:
            image_for_slide, image_source = find_image_for_slide(
                visuals_dir=visuals_dir,
                deck_title=deck_title,
                slide_title=slide_title if layout != "title" else deck_title,
                slide_number=slide_number,
                brain_output=brain_output,
                last_used_name=last_used_image_name
            )
        if image_for_slide:
            last_used_image_name = image_for_slide.name
            _mark_image_used(image_for_slide)

        stage_lower = clean(stage).lower()

        if image_source == "text_only":
            build_slide_text_only(slide, slide_title, body)
        elif layout == "title":
            add_base_background(slide)
            add_title_poster_image(slide, Path(POSTER_PATH) if POSTER_PATH else image_for_slide)
            add_top_rule(slide)
            add_title_text(slide, deck_title)
            place_text_by_stage(slide, stage, layout, body)
        else:
            # Full bleed — every non-title slide
            add_base_background(slide)
            add_full_bleed_image(slide, image_for_slide)
            add_top_rule(slide)
            if stage_lower != "closing":
                add_title_text(slide, slide_title.split("(")[0].strip())
            else:
                add_title_text(slide, deck_title)
            place_text_by_stage(slide, stage, layout, body)

        resolved_image_options = resolve_image_options_for_slide(
            visuals_dir=visuals_dir,
            slide_info=slide_info,
            image_for_slide=image_for_slide,
            image_source=image_source,
            slide_title=slide_title,
        )

        manifest.append({
            "slide_number": slide_number,
            "title": slide_title,
            "body": body,
            "layout": layout,
            "stage": stage,
            "image_path": str(image_for_slide) if image_for_slide else "",
            "image_name": image_for_slide.name if image_for_slide else "",
            "image_source": image_source,
            "image_query": slide_info.get("image_query", ""),
            "image_tags": slide_info.get("image_tags", []),
            "image_score": slide_info.get("image_score", 0),
            "image_options": resolved_image_options,
            "selected_option_id": resolved_image_options[0].get("option_id", "selected") if resolved_image_options else "",
        })

    out_path = next_output_path(output_dir)
    prs.save(str(out_path))
    manifest_path = output_dir / "latest_deck_manifest.json"
    manifest_path.write_text(json.dumps(manifest, indent=2), encoding="utf-8")
    print(f"📦 Deck manifest created: {manifest_path}")
    print(f"✅ Pitch deck created: {out_path}")
    return out_path


def parse_args():
    parser = argparse.ArgumentParser()
    parser.add_argument("slide_plan", nargs="?", help="Path to slide_plan.json")
    parser.add_argument("--project", help="Project/app directory containing slide_plan.json and visuals/")
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    slide_plan_path, visuals_dir, output_dir = resolve_paths(args)

    if not slide_plan_path.exists():
        print(f"Slide plan not found: {slide_plan_path}")
        raise SystemExit(1)

    build_presentation(slide_plan_path, visuals_dir, output_dir)


if __name__ == "__main__":
    main()
